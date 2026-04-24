from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


PPT_PATH = Path(r"F:\data_set_process\data_process\reports\group_meeting_model_progress_ppt_20260327\group_meeting_model_progress_20260327.pptx")
OUTPUT_PATH = Path(r"F:\data_set_process\data_process\reports\group_meeting_model_progress_ppt_20260327\group_meeting_model_progress_20260327_refined_from_markdown_v2.pptx")

COLORS = {
    "navy": RGBColor(26, 43, 68),
    "blue": RGBColor(62, 95, 138),
    "teal": RGBColor(69, 122, 120),
    "ink": RGBColor(37, 42, 48),
    "muted": RGBColor(97, 109, 126),
    "light_bg": RGBColor(245, 247, 250),
    "soft_line": RGBColor(214, 220, 228),
    "accent": RGBColor(186, 81, 64),
}


def add_text(slide, x, y, w, h, text, size=16, bold=False, color=None, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color or COLORS["ink"]


def add_bullets(slide, x, y, w, h, bullets, size=15.5, color=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for idx, line in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.bullet = True
        p.space_after = Pt(9)
        p.line_spacing = 1.1
        r = p.add_run()
        r.text = line
        r.font.name = "Microsoft YaHei"
        r.font.size = Pt(size)
        r.font.color.rgb = color or COLORS["ink"]


def add_panel(slide, x, y, w, h, title=None, fill_rgb=None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb or RGBColor(255, 255, 255)
    shape.line.color.rgb = COLORS["soft_line"]
    shape.line.width = Pt(1)
    if title:
        add_text(slide, x + Inches(0.18), y + Inches(0.1), w - Inches(0.36), Inches(0.26), title, size=11.5, bold=True, color=COLORS["blue"])


def add_metric_chip(slide, x, y, w, h, title, value, note, accent):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = COLORS["soft_line"]
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, Inches(0.08), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_text(slide, x + Inches(0.16), y + Inches(0.08), w - Inches(0.22), Inches(0.22), title, size=10.2, bold=True, color=COLORS["muted"])
    add_text(slide, x + Inches(0.16), y + Inches(0.29), w - Inches(0.22), Inches(0.26), value, size=17.5, bold=True, color=accent)
    add_text(slide, x + Inches(0.16), y + Inches(0.58), w - Inches(0.22), Inches(0.24), note, size=9.2, color=COLORS["muted"])


def add_info_card(slide, x, y, w, h, title, main_lines, note, accent, main_size=14.5):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = COLORS["soft_line"]
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, Inches(0.08), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    add_text(slide, x + Inches(0.16), y + Inches(0.08), w - Inches(0.22), Inches(0.22), title, size=10.2, bold=True, color=COLORS["muted"])
    add_text(slide, x + Inches(0.16), y + Inches(0.30), w - Inches(0.22), Inches(0.56), main_lines, size=main_size, bold=True, color=accent)
    add_text(slide, x + Inches(0.16), y + Inches(h) - Inches(0.23), w - Inches(0.22), Inches(0.16), note, size=8.8, color=COLORS["muted"])


def cover_content_area(slide):
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.4), Inches(0.45), Inches(12.55), Inches(6.78))
    rect.fill.solid()
    rect.fill.fore_color.rgb = COLORS["light_bg"]
    rect.line.fill.background()


def refine_slide1(slide):
    cover_content_area(slide)
    add_text(slide, Inches(0.72), Inches(0.90), Inches(10.9), Inches(0.56), "近期驾驶员短时控制行为预测模型进展汇报", size=26, bold=True)
    add_text(slide, Inches(0.74), Inches(1.42), Inches(10.5), Inches(0.38), "从基础轨迹预测到结构化行为建模的阶段性进展", size=15.5, color=COLORS["muted"])

    add_panel(slide, Inches(0.72), Inches(2.0), Inches(5.65), Inches(4.0), "这次汇报怎么讲")
    add_bullets(
        slide,
        Inches(0.95),
        Inches(2.42),
        Inches(5.1),
        Inches(2.95),
        [
            "重点不是简单罗列版本，而是围绕“当前遇到了什么问题、针对这些问题做了什么修改、效果怎么变化、现在还卡在哪里”来讲。",
            "从基础 baseline 到结构化版本，再到 deterministic conditioned v2，是一条连续演进路线，不是零散试验。",
            "当前已经能预测整体趋势，主问题收敛为：关键事件仍然对不准，而不是模型完全不会预测。",
        ],
        size=16.3,
    )

    add_panel(slide, Inches(6.7), Inches(2.0), Inches(5.9), Inches(4.0), "基本信息与重点")
    add_metric_chip(slide, Inches(6.95), Inches(2.48), Inches(1.72), Inches(0.98), "汇报人", "小郑", "按原稿保留", COLORS["teal"])
    add_metric_chip(slide, Inches(8.8), Inches(2.48), Inches(1.68), Inches(0.98), "时间", "2026.03", "组会汇报", COLORS["blue"])
    add_metric_chip(slide, Inches(10.6), Inches(2.48), Inches(1.7), Inches(0.98), "主推版本", "det v2", "当前主版本", COLORS["accent"])
    add_info_card(
        slide,
        Inches(6.95),
        Inches(3.72),
        Inches(5.35),
        Inches(1.08),
        "研究主题",
        "极限 / 高风险驾驶情境下驾驶员短时控制行为\n预测",
        "核心是“行为预测”，不是单纯未来曲线拟合",
        COLORS["blue"],
        main_size=12.5,
    )
    add_info_card(
        slide,
        Inches(6.95),
        Inches(4.98),
        Inches(5.35),
        Inches(1.02),
        "汇报目标",
        "让老师听明白：目前做到哪一步、\n为什么这样改、结果到底怎么样",
        "清楚、扎实、问题导向",
        COLORS["teal"],
        main_size=11.4,
    )
    add_text(slide, Inches(0.74), Inches(6.42), Inches(9.6), Inches(0.24), "汇报对象：导师 / 组会老师    风格：学术、简洁、清晰", size=10.2, color=COLORS["muted"])


def refine_slide2(slide):
    cover_content_area(slide)
    add_text(slide, Inches(0.58), Inches(0.62), Inches(8.5), Inches(0.38), "当前研究目标与核心问题", size=24, bold=True)
    add_text(slide, Inches(0.6), Inches(1.05), Inches(11.6), Inches(0.26), "任务已经从“预测未来 2 秒曲线”进一步收敛为“预测驾驶员短时控制行为”。", size=10.5, color=COLORS["muted"])

    add_panel(slide, Inches(0.72), Inches(1.72), Inches(5.9), Inches(4.9), "研究目标 / 当前主任务")
    add_bullets(
        slide,
        Inches(0.95),
        Inches(2.12),
        Inches(5.35),
        Inches(2.95),
        [
            "目标不是单纯拟合一条未来曲线，而是希望在危险 / 失稳情境下预测驾驶员接下来短时间内的控制行为。",
            "重点关注方向盘控制变化，并尽量让预测结果在趋势、关键动作和整体形态上接近真实驾驶员反应。",
            "输入：车辆历史状态、道路 / 场景相关实时信号；输出：未来短时间尺度内的方向盘控制轨迹，必要时辅以速度变化。",
        ],
        size=14.9,
    )

    add_panel(slide, Inches(6.78), Inches(1.72), Inches(5.82), Inches(4.9), "当前识别出的主要问题")
    add_bullets(
        slide,
        Inches(7.02),
        Inches(2.12),
        Inches(5.28),
        Inches(2.9),
        [
            "基础模型前段有时能够跟上一部分趋势，但后段容易偏离真实轨迹。",
            "关键事件预测不准，主要体现在起转时机、主峰位置、反打 / 纠偏幅度等方面。",
            "在 interaction 等高歧义场景下，单一输出容易形成折中结果，不够像真实驾驶员行为。",
        ],
        size=14.9,
    )

    strip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.92), Inches(6.08), Inches(11.55), Inches(0.60))
    strip.fill.solid()
    strip.fill.fore_color.rgb = RGBColor(236, 242, 248)
    strip.line.fill.background()
    add_text(slide, Inches(1.16), Inches(6.25), Inches(11.0), Inches(0.18), "讲述重点：目前最主要的问题，不只是整体跟踪不稳，更重要的是关键动作不准；interaction 场景还提示这个问题可能天然存在多种合理未来。", size=11.2, bold=True, color=COLORS["navy"])


def main():
    prs = Presentation(str(PPT_PATH))
    refine_slide1(prs.slides[0])
    refine_slide2(prs.slides[1])
    prs.save(str(OUTPUT_PATH))


if __name__ == "__main__":
    main()
