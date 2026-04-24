from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SRC_PPT = Path(r"F:\data_set_process\data_process\reports\group_meeting_model_progress_ppt_20260327\group_meeting_model_progress_20260327_refined_from_markdown_v2.pptx")
OUT_PPT = Path(r"F:\data_set_process\data_process\reports\group_meeting_model_progress_ppt_20260327\group_meeting_model_progress_20260327_refined_from_markdown_v3_same_pool_baseline.pptx")
NEW_IMG = Path(r"F:\data_set_process\data_process\reports\fair_baseline_same_pool_check_20260328\fair_same_pool_representative_samples_overview.png")

LIGHT_BG = RGBColor(245, 247, 250)
SOFT_LINE = RGBColor(214, 220, 228)
MUTED = RGBColor(97, 109, 126)
NAVY = RGBColor(26, 43, 68)


def add_text(slide, x, y, w, h, text, size=10.5, bold=False, color=None, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Microsoft YaHei"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color or NAVY


def cover(slide, x, y, w, h):
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, w, h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = LIGHT_BG
    rect.line.fill.background()


def add_frame(slide, x, y, w, h):
    rect = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor(255, 255, 255)
    rect.line.color.rgb = SOFT_LINE
    rect.line.width = Pt(1)


def main():
    prs = Presentation(str(SRC_PPT))
    slide = prs.slides[2]

    # Replace the right-side old baseline illustration with the same-pool formal comparison overview.
    cover(slide, Inches(5.0), Inches(1.45), Inches(7.9), Inches(5.35))
    add_frame(slide, Inches(5.05), Inches(1.62), Inches(7.55), Inches(4.78))
    slide.shapes.add_picture(str(NEW_IMG), Inches(5.22), Inches(1.82), Inches(7.18), Inches(4.18))

    # Replace the old caption with a same-pool caption.
    cover(slide, Inches(5.05), Inches(6.18), Inches(7.55), Inches(0.38))
    add_text(
        slide,
        Inches(5.18),
        Inches(6.28),
        Inches(7.15),
        Inches(0.18),
        "图：基础 baseline 在同一份 6238 样本池（train 4797 / val 692 / test 749）下的代表性预测-真实方向盘转角曲线对照。",
        size=9.8,
        color=MUTED,
    )

    # Add a small fairness note near the image block so the slide is self-explanatory.
    cover(slide, Inches(5.08), Inches(1.5), Inches(4.3), Inches(0.2))
    add_text(
        slide,
        Inches(5.15),
        Inches(1.52),
        Inches(4.2),
        Inches(0.16),
        "已替换为与当前主版本相同 sample pool 的 baseline 对照图",
        size=9.5,
        bold=True,
        color=NAVY,
    )

    prs.save(str(OUT_PPT))


if __name__ == "__main__":
    main()
