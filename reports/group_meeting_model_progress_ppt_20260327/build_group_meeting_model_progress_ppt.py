from __future__ import annotations

from pathlib import Path

import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"F:\data_set_process\data_process")
WORKDIR = ROOT / "reports" / "group_meeting_model_progress_ppt_20260327"
ASSET_DIR = WORKDIR / "assets"
OUTPUT_PPTX = WORKDIR / "group_meeting_model_progress_20260327.pptx"


COLORS = {
    "navy": RGBColor(26, 43, 68),
    "blue": RGBColor(62, 95, 138),
    "teal": RGBColor(69, 122, 120),
    "ink": RGBColor(37, 42, 48),
    "muted": RGBColor(97, 109, 126),
    "light_bg": RGBColor(245, 247, 250),
    "soft_line": RGBColor(214, 220, 228),
    "accent": RGBColor(186, 81, 64),
    "good": RGBColor(42, 110, 93),
}

FONT_CN = "Microsoft YaHei"
FONT_EN = "Aptos"


def ensure_dirs() -> None:
    WORKDIR.mkdir(parents=True, exist_ok=True)
    ASSET_DIR.mkdir(parents=True, exist_ok=True)


def create_selection_rule_plot(output_path: Path) -> None:
    labels = [
        "Primary RMSE\n(lower better)",
        "Tail RMSE\n(lower better)",
        "Boundary Shift\n(lower better)",
        "Turning Count Err\n(lower better)",
    ]
    legacy = [0.3665, 0.3691, 0.7358, 1.4597]
    corrected = [0.3829, 0.3631, 0.6340, 1.4355]

    plt.rcParams["font.sans-serif"] = ["Microsoft YaHei", "SimHei", "DejaVu Sans"]
    plt.rcParams["axes.unicode_minus"] = False

    fig, ax = plt.subplots(figsize=(10.8, 5.6), dpi=180)
    fig.patch.set_facecolor("#f7f8fb")
    ax.set_facecolor("#f7f8fb")

    x = range(len(labels))
    width = 0.33
    ax.bar([i - width / 2 for i in x], legacy, width=width, color="#c7d2e0", label="Legacy rule (epoch 3)")
    ax.bar([i + width / 2 for i in x], corrected, width=width, color="#406d93", label="Corrected rule (epoch 1)")

    for idx, (a, b) in enumerate(zip(legacy, corrected)):
        ax.text(idx - width / 2, a + 0.03, f"{a:.3f}", ha="center", va="bottom", fontsize=10, color="#334155")
        ax.text(idx + width / 2, b + 0.03, f"{b:.3f}", ha="center", va="bottom", fontsize=10, color="#1f2937")
        marker = "Improved" if b < a else "Trade-off"
        color = "#2f6f5f" if b < a else "#b45309"
        ax.text(idx, max(a, b) + 0.16, marker, ha="center", va="bottom", fontsize=10, color=color, fontweight="bold")

    ax.set_xticks(list(x))
    ax.set_xticklabels(labels, fontsize=11)
    ax.set_ylim(0, 1.95)
    ax.grid(axis="y", linestyle="--", alpha=0.28)
    ax.spines[["top", "right", "left"]].set_visible(False)
    ax.spines["bottom"].set_color("#cbd5e1")
    ax.tick_params(axis="y", length=0)
    ax.set_yticklabels([])
    ax.legend(loc="upper right", frameon=False, fontsize=11)
    ax.set_title(
        "Selection Rule Correction: legacy primary-RMSE checkpoint vs structure-aware checkpoint",
        fontsize=15,
        fontweight="bold",
        color="#1f2937",
        pad=16,
    )
    ax.text(
        0.0,
        1.02,
        "Legacy rule optimizes a 0-1.5 s primary metric; corrected rule keeps the checkpoint that is better on tail/structure behavior.",
        transform=ax.transAxes,
        fontsize=11,
        color="#475569",
    )

    plt.tight_layout()
    fig.savefig(output_path, bbox_inches="tight")
    plt.close(fig)


def setup_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def add_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = COLORS["light_bg"]

    header = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.42))
    header.fill.solid()
    header.fill.fore_color.rgb = COLORS["navy"]
    header.line.fill.background()

    footer = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(7.28), Inches(13.333), Inches(0.22))
    footer.fill.solid()
    footer.fill.fore_color.rgb = RGBColor(230, 234, 240)
    footer.line.fill.background()


def add_title(slide, title: str, subtitle: str | None = None, page_no: int | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.58), Inches(9.8), Inches(0.58))
    p = title_box.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = FONT_CN
    r.font.size = Pt(24)
    r.font.bold = True
    r.font.color.rgb = COLORS["ink"]

    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.58), Inches(1.08), Inches(10.8), Inches(0.42))
        p = subtitle_box.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = FONT_CN
        r.font.size = Pt(10.5)
        r.font.color.rgb = COLORS["muted"]

    if page_no is not None:
        num_box = slide.shapes.add_textbox(Inches(12.35), Inches(0.62), Inches(0.45), Inches(0.28))
        p = num_box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = str(page_no)
        r.font.name = FONT_EN
        r.font.size = Pt(12)
        r.font.bold = True
        r.font.color.rgb = COLORS["muted"]


def add_panel(slide, x, y, w, h, title: str | None = None, fill_rgb: RGBColor | None = None):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb or RGBColor(255, 255, 255)
    shape.line.color.rgb = COLORS["soft_line"]
    shape.line.width = Pt(1)

    if title:
        tb = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.1), w - Inches(0.36), Inches(0.26))
        p = tb.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.name = FONT_CN
        r.font.size = Pt(11.5)
        r.font.bold = True
        r.font.color.rgb = COLORS["blue"]
    return shape


def add_bullets(slide, x, y, w, h, bullets: list[str], font_size: int = 16, color: RGBColor | None = None) -> None:
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(3)
    tf.margin_right = Pt(3)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for idx, line in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(9)
        p.line_spacing = 1.12
        p.bullet = True
        r = p.add_run()
        r.text = line
        r.font.name = FONT_CN
        r.font.size = Pt(font_size)
        r.font.color.rgb = color or COLORS["ink"]


def add_text(slide, x, y, w, h, text: str, size=16, bold=False, color=None, align=PP_ALIGN.LEFT) -> None:
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = FONT_CN
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color or COLORS["ink"]


def add_metric_chip(slide, x, y, w, h, title: str, value: str, note: str, accent: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = COLORS["soft_line"]

    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x, y, Inches(0.08), h)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    add_text(slide, x + Inches(0.16), y + Inches(0.10), w - Inches(0.24), Inches(0.22), title, size=10.5, bold=True, color=COLORS["muted"])
    add_text(slide, x + Inches(0.16), y + Inches(0.31), w - Inches(0.24), Inches(0.28), value, size=18, bold=True, color=accent)
    add_text(slide, x + Inches(0.16), y + Inches(0.62), w - Inches(0.24), Inches(0.28), note, size=9.5, color=COLORS["muted"])


def add_picture(slide, image_path: Path, x, y, w, h) -> None:
    slide.shapes.add_picture(str(image_path), x, y, w, h)


def add_caption(slide, x, y, w, text: str) -> None:
    add_text(slide, x, y, w, Inches(0.22), text, size=9.5, color=COLORS["muted"])


def build_deck() -> None:
    ensure_dirs()

    selection_plot = ASSET_DIR / "selection_rule_comparison.png"
    create_selection_rule_plot(selection_plot)

    prs = setup_presentation()

    img_baseline = ROOT / "reports" / "single_output_reinforcement_d3_gpu_v4" / "figures" / "old_style_sampling_current_model" / "pred_vs_gt_examples_old_sampling_seed2026_steer_only.png"
    img_structured_v1 = ROOT / "reports" / "event_plus_conditioned_trajectory_baseline_20260326" / "task_D_formal_run" / "figures" / "representative_samples_overview.png"
    img_v2 = ROOT / "reports" / "v3_selection_conditioned_interaction_pilot_20260327" / "task_2_conditioned_v2" / "formal_eval" / "figures" / "representative_samples_overview.png"
    img_multihyp = ROOT / "reports" / "v3_selection_conditioned_interaction_pilot_20260327" / "task_3_interaction_multihyp" / "formal_eval" / "figures" / "representative_samples_overview.png"

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_text(slide, Inches(0.72), Inches(0.95), Inches(8.9), Inches(0.52), "模型进展组会汇报", size=28, bold=True, color=COLORS["ink"])
    add_text(slide, Inches(0.74), Inches(1.55), Inches(9.6), Inches(0.48), "核心目标：把“为什么连续改版本、每次解决了什么、现在到底卡在哪里”讲清楚。", size=15, color=COLORS["muted"])
    add_panel(slide, Inches(0.72), Inches(2.05), Inches(5.65), Inches(3.9), "本次汇报想回答的三个问题")
    add_bullets(
        slide,
        Inches(0.95),
        Inches(2.52),
        Inches(5.1),
        Inches(3.0),
        [
            "我的研究主线不是零散试验，而是围绕“2 s 方向盘转角预测中的关键事件对齐”持续收敛。",
            "从 baseline 到结构化 conditioned，再到 deterministic conditioned v2，每次修改都对应一个具体问题。",
            "当前已经能稳定预测整体趋势，主问题收敛为：关键事件仍然对不准，而不是模型完全不会预测。",
        ],
        font_size=17,
    )
    add_panel(slide, Inches(6.7), Inches(2.05), Inches(5.9), Inches(3.9), "这次重点")
    add_metric_chip(slide, Inches(6.95), Inches(2.55), Inches(2.9), Inches(1.0), "当前主推版本", "deterministic conditioned v2", "当前最平衡、最值得继续深挖的主线版本", COLORS["blue"])
    add_metric_chip(slide, Inches(9.95), Inches(2.55), Inches(2.35), Inches(1.0), "关键修正", "selection rule", "直接影响 best checkpoint 是否选对", COLORS["accent"])
    add_metric_chip(slide, Inches(6.95), Inches(3.82), Inches(5.35), Inches(1.0), "当前结论", "问题已收敛到“关键事件对齐”", "interaction-only multi-hypothesis pilot 证明方向值得继续，但还不是成熟方案", COLORS["teal"])
    add_text(slide, Inches(0.74), Inches(6.35), Inches(8.9), Inches(0.28), "汇报对象：导师 / 组会老师    口径：问题导向、结论优先、不过度包装", size=10.5, color=COLORS["muted"])

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "研究主线与版本演进", "按“发现问题 -> 修改机制 -> 修正评估 -> 收敛主版本”的逻辑看，而不是把实验当成孤立试错。", 2)
    stages = [
        ("1. 基础 baseline", "先确认模型能否学到 2 s 整体转角趋势", "整体趋势能学到，但 tail / turning / reversal 对不准"),
        ("2. 第一代结构化版本", "把事件信息显式注入轨迹预测", "局部结构开始改善，但收益不稳定"),
        ("3. selection rule 修正", "避免用 0-1.5 s 主 RMSE 选错 checkpoint", "评估目标开始和研究目标一致"),
        ("4. deterministic conditioned v2", "强化事件-轨迹耦合并让条件更可控", "目前最平衡、最适合作为主版本"),
        ("5. interaction-only multihyp pilot", "验证交互场景是否真的需要多假设", "oracle 有收益，说明方向值得继续"),
    ]
    base_x = 0.75
    width = 2.35
    y = 1.85
    for i, (title, why, outcome) in enumerate(stages):
        x = Inches(base_x + i * 2.45)
        add_panel(slide, x, Inches(y), Inches(width), Inches(3.9), None)
        circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, x + Inches(0.14), Inches(y) + Inches(0.18), Inches(0.42), Inches(0.42))
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLORS["blue"] if i < 3 else COLORS["teal"]
        circle.line.fill.background()
        add_text(slide, x + Inches(0.67), Inches(y) + Inches(0.15), Inches(1.45), Inches(0.38), title, size=12, bold=True)
        add_text(slide, x + Inches(0.16), Inches(y) + Inches(0.72), Inches(2.0), Inches(1.02), why, size=11, color=COLORS["muted"])
        add_text(slide, x + Inches(0.16), Inches(y) + Inches(2.05), Inches(2.02), Inches(1.45), outcome, size=11.2, bold=(i == 3), color=COLORS["ink"])
        if i < len(stages) - 1:
            arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, x + Inches(2.22), Inches(y) + Inches(1.68), Inches(0.22), Inches(0.28))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = COLORS["soft_line"]
            arrow.line.fill.background()
    add_panel(slide, Inches(0.84), Inches(6.0), Inches(11.8), Inches(0.78), None, RGBColor(236, 242, 248))
    add_text(
        slide,
        Inches(1.05),
        Inches(6.18),
        Inches(11.3),
        Inches(0.34),
        "主线结论：模型已经从“学不到”过渡到“能学到，但关键事件仍偏移”；因此后续工作不再是盲目堆模型，而是继续缩小事件锚点误差。",
        size=13,
        bold=True,
        color=COLORS["navy"],
    )

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "基础 Baseline：已学到整体趋势，但关键事件仍然对不准", "代表图使用项目现有的预测-真实方向盘转角曲线对比。", 3)
    add_panel(slide, Inches(0.62), Inches(1.58), Inches(4.2), Inches(4.95), "我先用 baseline 想回答什么问题")
    add_bullets(
        slide,
        Inches(0.84),
        Inches(2.0),
        Inches(3.76),
        Inches(2.45),
        [
            "先确认纯 2 s 监督下，模型是否至少能学到方向盘转角的整体变化趋势。",
            "如果 baseline 连整体都学不到，后面的结构化设计就没有意义。",
            "因此 baseline 的作用不是追求最终最优，而是给后续修改提供参照系。",
        ],
        font_size=15.5,
    )
    add_metric_chip(slide, Inches(0.85), Inches(4.82), Inches(1.7), Inches(0.95), "2 s RMSE", "0.381", "同一评估协议下的 unconditional baseline", COLORS["blue"])
    add_metric_chip(slide, Inches(2.68), Inches(4.82), Inches(1.95), Inches(0.95), "Tail RMSE", "0.398", "1.5-2.0 s 尾段误差偏大", COLORS["accent"])
    add_metric_chip(slide, Inches(0.85), Inches(5.82), Inches(1.7), Inches(0.95), "Tail Trend Corr", "0.066", "尾段趋势相关性很弱", COLORS["accent"])
    add_metric_chip(slide, Inches(2.68), Inches(5.82), Inches(1.95), Inches(0.95), "Turning Count Err", "1.77", "转折次数经常对不上", COLORS["accent"])
    add_picture(slide, img_baseline, Inches(5.12), Inches(1.72), Inches(7.45), Inches(4.58))
    add_caption(slide, Inches(5.18), Inches(6.34), Inches(7.2), "图：基础 baseline 的代表性预测-真实方向盘转角曲线对比（项目现有结果图）")

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "第一代结构化版本：把事件信息显式注入轨迹预测", "这一版的意义在于把问题从“纯回归”转成“事件引导的轨迹生成”。", 4)
    add_panel(slide, Inches(0.62), Inches(1.55), Inches(4.3), Inches(5.05), "这一版具体改了什么")
    add_bullets(
        slide,
        Inches(0.84),
        Inches(1.98),
        Inches(3.88),
        Inches(2.3),
        [
            "在 baseline 上加入事件条件，让模型不只看历史轨迹，还显式接收 turn onset / peak / reversal 等结构信息。",
            "目的不是简单加特征，而是希望模型在关键节点附近更敢于拐、该回的时候能回。",
            "这一版已经出现明显的结构收益，但收益并不稳定。不同切片上表现分化较大。",
        ],
        font_size=15,
    )
    add_metric_chip(slide, Inches(0.84), Inches(4.68), Inches(1.86), Inches(0.95), "Primary Tail RMSE", "0.390 -> 0.349", "primary 子集尾段误差明显下降", COLORS["good"])
    add_metric_chip(slide, Inches(2.82), Inches(4.68), Inches(1.86), Inches(0.95), "Interaction Tail RMSE", "0.495 -> 0.436", "interaction 切片也有收益", COLORS["good"])
    add_metric_chip(slide, Inches(0.84), Inches(5.72), Inches(1.86), Inches(0.95), "Overall 2 s RMSE", "0.381 -> 0.388", "全局误差未必同步变好", COLORS["accent"])
    add_metric_chip(slide, Inches(2.82), Inches(5.72), Inches(1.86), Inches(0.95), "Tail Trend Corr", "0.066 -> 0.026", "事件对齐仍然不稳定", COLORS["accent"])
    add_picture(slide, img_structured_v1, Inches(5.15), Inches(1.72), Inches(7.4), Inches(4.58))
    add_caption(slide, Inches(5.2), Inches(6.33), Inches(7.1), "图：第一代结构化版本代表图。可以看到部分转折结构更像，但稳定性仍不够。")

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "selection rule 修正：先把 best checkpoint 选对", "这一步很关键，因为如果 checkpoint 选错，后面所有版本判断都会失真。", 5)
    add_panel(slide, Inches(0.62), Inches(1.55), Inches(4.25), Inches(5.0), "为什么必须改 selection rule")
    add_bullets(
        slide,
        Inches(0.84),
        Inches(1.98),
        Inches(3.78),
        Inches(2.5),
        [
            "训练审计发现：原来的 best checkpoint 由 val sample-weighted overall primary steer RMSE 决定，本质上更偏 0-1.5 s 指标。",
            "但我的研究目标已经转向 tail / turning / boundary continuity，所以旧规则会把“看起来 RMSE 更低、但结构更差”的模型选出来。",
            "修正后改为 structure-aware / active 的选择逻辑，让 checkpoint 选择和研究目标一致。",
        ],
        font_size=15,
    )
    add_metric_chip(slide, Inches(0.84), Inches(4.92), Inches(1.8), Inches(0.95), "Legacy picks", "epoch 3", "主 RMSE 更低，但不是我要的结构最优", COLORS["accent"])
    add_metric_chip(slide, Inches(2.78), Inches(4.92), Inches(1.8), Inches(0.95), "Corrected picks", "epoch 1", "尾段与边界结构更合理", COLORS["good"])
    add_metric_chip(slide, Inches(0.84), Inches(5.93), Inches(1.8), Inches(0.95), "Boundary Shift", "0.736 -> 0.634", "修正规则后明显下降", COLORS["good"])
    add_metric_chip(slide, Inches(2.78), Inches(5.93), Inches(1.8), Inches(0.95), "Tail RMSE", "0.369 -> 0.363", "虽然不是巨大跃升，但方向正确", COLORS["good"])
    add_picture(slide, selection_plot, Inches(5.08), Inches(1.8), Inches(7.48), Inches(4.72))
    add_caption(slide, Inches(5.15), Inches(6.34), Inches(7.2), "图：selection rule 修正前后 best checkpoint 的测试集表现变化。Primary RMSE 略有 trade-off，但 tail / boundary / turning 更符合当前目标。")

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "deterministic conditioned v2：当前最重要的主版本", "这一页重点回答：改了什么、为什么这样改、比前一版和 baseline 好在哪里、还差什么。", 6)
    add_panel(slide, Inches(0.55), Inches(1.48), Inches(4.45), Inches(5.18), "为什么说它是当前主推版本")
    add_bullets(
        slide,
        Inches(0.78),
        Inches(1.9),
        Inches(3.98),
        Inches(2.45),
        [
            "在第一代结构化版本基础上，改成 `structured_v2` conditioning，并加入更强的 deterministic event-to-trajectory coupling。",
            "核心配置包括：`structure_width=0.065`、`gate_temperature=0.04`、`event_residual_scale=1.2`，同时配合修正后的 structure-aware selection。",
            "目标是让事件条件不再“软注入”，而是更稳定地约束关键转折位置和尾段走势。",
        ],
        font_size=14.6,
    )
    add_metric_chip(slide, Inches(0.78), Inches(4.76), Inches(1.92), Inches(0.94), "Overall 2 s RMSE", "0.381 -> 0.377", "相对 baseline 小幅下降", COLORS["good"])
    add_metric_chip(slide, Inches(2.84), Inches(4.76), Inches(1.92), Inches(0.94), "Tail RMSE", "0.398 -> 0.376", "尾段误差更稳定下降", COLORS["good"])
    add_metric_chip(slide, Inches(0.78), Inches(5.78), Inches(1.92), Inches(0.94), "Interaction Tail RMSE", "0.495 -> 0.421", "交互切片收益更明显", COLORS["good"])
    add_metric_chip(slide, Inches(2.84), Inches(5.78), Inches(1.92), Inches(0.94), "Turning Count Err", "1.77 -> 1.54", "转折数误差明显下降", COLORS["good"])
    add_picture(slide, img_v2, Inches(5.16), Inches(1.72), Inches(7.35), Inches(3.62))
    add_caption(slide, Inches(5.2), Inches(5.42), Inches(7.15), "图：deterministic conditioned v2 的代表结果图（预测 vs 真实方向盘转角曲线）")
    add_panel(slide, Inches(5.15), Inches(5.72), Inches(7.38), Inches(0.96), None, RGBColor(234, 242, 241))
    add_text(
        slide,
        Inches(5.36),
        Inches(5.92),
        Inches(6.95),
        Inches(0.3),
        "当前判断：它比 baseline 和上一版更像“真正围绕关键事件在预测”，但主瓶颈仍是 onset / reversal / peak 的时间对齐不够准，interaction 场景尤其明显。",
        size=11.2,
        bold=True,
        color=COLORS["navy"],
    )

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "interaction-only multi-hypothesis pilot：方向值得继续，但现在还不是成熟方案", "只在交互场景上做 pilot，目的是验证“多假设”到底有没有必要。", 7)
    add_panel(slide, Inches(0.62), Inches(1.55), Inches(4.3), Inches(5.04), "pilot 想验证什么")
    add_bullets(
        slide,
        Inches(0.84),
        Inches(1.98),
        Inches(3.82),
        Inches(2.38),
        [
            "交互场景本身存在未来不确定性，所以我没有直接全量上多模态，而是先做 interaction-only multi-hypothesis pilot。",
            "如果 oracle 明显优于 deterministic v2，说明“多假设本身有价值”；如果 top-1 选不出来，则问题在 selector，而不一定在 hypothesis 生成。",
            "这个 pilot 的定位是“验证方向”，不是直接替代当前主版本。",
        ],
        font_size=14.8,
    )
    add_metric_chip(slide, Inches(0.84), Inches(4.75), Inches(1.78), Inches(0.95), "Det v2", "Tail RMSE 0.344", "当前稳定基线", COLORS["blue"])
    add_metric_chip(slide, Inches(2.72), Inches(4.75), Inches(1.92), Inches(0.95), "Top-1 multihyp", "Tail RMSE 0.387", "当前 selector 还不成熟", COLORS["accent"])
    add_metric_chip(slide, Inches(0.84), Inches(5.78), Inches(1.78), Inches(0.95), "Oracle multihyp", "Tail RMSE 0.256", "说明多假设本身有潜力", COLORS["good"])
    add_metric_chip(slide, Inches(2.72), Inches(5.78), Inches(1.92), Inches(0.95), "Overall judgment", "值得继续", "先解决 hypothesis selection", COLORS["teal"])
    add_picture(slide, img_multihyp, Inches(5.08), Inches(1.72), Inches(7.45), Inches(4.58))
    add_caption(slide, Inches(5.16), Inches(6.33), Inches(7.15), "图：interaction multi-hypothesis pilot 代表图。oracle 收益存在，但 top-1 选择仍不稳定。")

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "总结与下一步工作", "最后一页的口径是：我已经做到哪一步、主要瓶颈在哪里、下一步准备怎么做。", 8)
    add_panel(slide, Inches(0.6), Inches(1.55), Inches(3.95), Inches(4.95), "已经取得的进展")
    add_bullets(
        slide,
        Inches(0.84),
        Inches(1.98),
        Inches(3.46),
        Inches(2.58),
        [
            "已经建立起从 baseline 到 deterministic conditioned v2 的连续演进主线，而不是停留在零散试验。",
            "已经把“结构化建模”和“selection rule 修正”两件事拆开讲清楚，并验证它们分别带来的收益。",
            "当前主推版本 deterministic conditioned v2 在 tail、turning、interaction 切片上都比 baseline 更好。",
        ],
        font_size=14.7,
    )
    add_panel(slide, Inches(4.72), Inches(1.55), Inches(3.95), Inches(4.95), "当前主要瓶颈")
    add_bullets(
        slide,
        Inches(4.96),
        Inches(1.98),
        Inches(3.44),
        Inches(2.58),
        [
            "问题已经收敛为：关键事件仍然对不准，尤其是 onset / reversal / peak 的时间锚点误差。",
            "也就是说，模型不是“整体完全不会预测”，而是“已经会预测，但关键节点不够准”。",
            "interaction 场景仍然是最难的部分，说明未来不确定性和 selector 设计都是核心难点。",
        ],
        font_size=14.7,
    )
    add_panel(slide, Inches(8.84), Inches(1.55), Inches(3.88), Inches(4.95), "下一步准备怎么做")
    add_bullets(
        slide,
        Inches(9.08),
        Inches(1.98),
        Inches(3.38),
        Inches(2.58),
        [
            "继续围绕 deterministic conditioned v2 做事件锚点精修，而不是立刻换主线。",
            "补强 turning / reversal / peak timing 的显式约束，让“关键事件对齐”进入 loss 和 selection。",
            "把 multi-hypothesis 先限定在 interaction slice，重点先解决 hypothesis selector，再决定是否扩展。",
        ],
        font_size=14.4,
    )
    add_panel(slide, Inches(0.72), Inches(6.02), Inches(12.0), Inches(0.72), None, RGBColor(233, 239, 245))
    add_text(
        slide,
        Inches(0.96),
        Inches(6.2),
        Inches(11.55),
        Inches(0.28),
        "一句话总结：我现在已经把问题从“模型能不能预测”收敛到了“关键事件如何更准确对齐”；因此下一步不是盲目加复杂度，而是继续沿着 deterministic conditioned v2 把事件锚点做准。",
        size=12.2,
        bold=True,
        color=COLORS["navy"],
    )

    prs.save(str(OUTPUT_PPTX))


if __name__ == "__main__":
    build_deck()
